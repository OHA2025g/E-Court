"""PowerPoint dashboard export with KPI, trend, state RAG grid, and HC rankings."""
import io
from typing import Callable, Optional

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

from dashboard_agg import (
    compute_dashboard_by_hc,
    compute_dashboard_summary,
    compute_states_rag,
    compute_trend_with_milestones,
)

RAG_HEX = {
    "GREEN": "#10B981",
    "AMBER": "#F59E0B",
    "RED": "#EF4444",
    "NA": "#94A3B8",
}


def _hex_rgb(h: str) -> tuple:
    h = h.lstrip("#")
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _fmt_pct(v) -> str:
    if v is None:
        return "—"
    return f"{v:.1f}%"


def render_trend_png(trend: dict, width: int = 960, height: int = 480) -> bytes:
    """Line chart of physical achievement % by reporting period."""
    periods = trend.get("periods") or []
    if not periods:
        img = Image.new("RGB", (width, height), "white")
        d = ImageDraw.Draw(img)
        d.text((20, 20), "No trend data available", fill="#334155")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    values = [p.get("phys_percent") or 0 for p in periods]
    labels = [p.get("period", "")[-5:] for p in periods]
    margin = 60
    plot_w = width - margin * 2
    plot_h = height - margin * 2
    vmax = max(values) if values else 100
    vmax = max(vmax, 1)

    img = Image.new("RGB", (width, height), "white")
    d = ImageDraw.Draw(img)
    d.text((margin, 12), "Physical achievement trend (approved data)", fill="#0A1128")

    pts = []
    n = len(values)
    for i, v in enumerate(values):
        x = margin + (plot_w * i / max(n - 1, 1))
        y = margin + plot_h - (plot_h * v / vmax)
        pts.append((x, y))

    if len(pts) >= 2:
        d.line(pts, fill=_hex_rgb(RAG_HEX["GREEN"]), width=3)
    for (x, y), lbl in zip(pts, labels):
        d.ellipse((x - 4, y - 4, x + 4, y + 4), fill=_hex_rgb(RAG_HEX["GREEN"]))
        d.text((x - 12, height - margin + 8), lbl, fill="#64748B")

    d.line((margin, margin + plot_h, margin + plot_w, margin + plot_h), fill="#CBD5E1")
    d.line((margin, margin, margin, margin + plot_h), fill="#CBD5E1")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def render_states_rag_png(states: dict, width: int = 960, height: int = 540) -> bytes:
    """Grid of state/UT cells coloured by parent HC physical RAG."""
    img = Image.new("RGB", (width, height), "white")
    d = ImageDraw.Draw(img)
    d.text((16, 10), "India — physical RAG by state/UT jurisdiction", fill="#0A1128")

    items = sorted(states.items())
    if not items:
        d.text((20, 40), "No state data", fill="#64748B")
    else:
        cols = 6
        cell_w = (width - 32) // cols
        cell_h = 36
        for idx, (state, info) in enumerate(items):
            row, col = divmod(idx, cols)
            x = 16 + col * cell_w
            y = 40 + row * cell_h
            rag = info.get("rag", "NA") if info.get("in_scope", True) else "NA"
            fill = _hex_rgb(RAG_HEX.get(rag, RAG_HEX["NA"]))
            d.rectangle((x, y, x + cell_w - 4, y + cell_h - 6), fill=fill, outline="#E2E8F0")
            label = state[:14] + ("…" if len(state) > 14 else "")
            d.text((x + 4, y + 8), label, fill="#0A1128")

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _add_bullet_slide(prs, title: str, lines: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = box.text_frame
    tf.text = lines[0] if lines else ""
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.level = 0


def _add_image_slide(prs, title: str, png_bytes: bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.4), Inches(1.3), width=Inches(9.2))


async def build_dashboard_pptx(
    db,
    scope_filter_fn: Callable,
    compute_rag_fn: Callable,
    safe_div_fn: Callable,
    user: dict,
    reporting_period: Optional[str],
    extra_match: Optional[dict],
    state_to_hc: dict,
) -> bytes:
    summary = await compute_dashboard_summary(
        db, scope_filter_fn, compute_rag_fn, safe_div_fn, user, reporting_period, extra_match,
    )
    trend = await compute_trend_with_milestones(
        db, scope_filter_fn, safe_div_fn, user, extra_match,
    )
    states = await compute_states_rag(
        db, state_to_hc, scope_filter_fn, compute_rag_fn, user, reporting_period, "physical", extra_match,
    )
    hc_rows = await compute_dashboard_by_hc(
        db, scope_filter_fn, safe_div_fn, user, reporting_period, extra_match,
    )

    phys = summary.get("physical", {})
    fin = summary.get("financial", {})
    rag_p = summary.get("rag_physical", {})
    rag_f = summary.get("rag_financial", {})
    label = reporting_period or "Latest"

    ranked = [r for r in hc_rows if r.get("phys_percent") is not None]
    ranked.sort(key=lambda x: x["phys_percent"], reverse=True)
    top5 = ranked[:5]
    bottom5 = list(reversed(ranked[-5:])) if len(ranked) >= 5 else list(reversed(ranked))

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"eCourts PMIS Dashboard — {label}"
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = box.text_frame
    tf.text = "National KPI Summary (approved data only)"
    kpi_lines = [
        f"Physical achievement: {_fmt_pct(phys.get('percent'))}",
        f"Financial utilisation: {_fmt_pct(fin.get('utilisation_percent'))}",
        f"Physical RAG — Green: {rag_p.get('GREEN', 0)}, Amber: {rag_p.get('AMBER', 0)}, Red: {rag_p.get('RED', 0)}",
        f"Financial RAG — Green: {rag_f.get('GREEN', 0)}, Amber: {rag_f.get('AMBER', 0)}, Red: {rag_f.get('RED', 0)}",
        f"Outcome KPI rows tracked: {summary.get('outcome', {}).get('kpi_count', 0)}",
    ]
    for line in kpi_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    _add_image_slide(prs, "National trend — physical achievement", render_trend_png(trend))
    _add_image_slide(prs, "India choropleth — physical RAG by state", render_states_rag_png(states))

    top_lines = [f"{r['high_court']}: {_fmt_pct(r['phys_percent'])}" for r in top5]
    bot_lines = [f"{r['high_court']}: {_fmt_pct(r['phys_percent'])}" for r in bottom5]
    _add_bullet_slide(
        prs,
        f"Top & bottom High Courts — physical ({label})",
        ["Top performers:"] + top_lines + [""] + ["Needs attention:"] + bot_lines,
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
